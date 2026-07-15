import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor


def create_ppt(
    presentation_text: str,
    output_path: str,
):
    """
    Convert AI generated presentation text
    into a professional PowerPoint presentation.
    """

    prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = re.split(
        r"###\s*Slide\s*\d+",
        presentation_text,
    )

    slide_number = 1

    for slide in slides:

        slide = slide.strip()

        if not slide:
            continue

        # ----------------------------------------
        # Layout
        # ----------------------------------------
        if slide_number == 1:
            layout = prs.slide_layouts[0]
        else:
            layout = prs.slide_layouts[1]

        ppt_slide = prs.slides.add_slide(layout)

        # ----------------------------------------
        # Blue Header
        # ----------------------------------------
        banner = ppt_slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            0,
            0,
            prs.slide_width,
            Inches(0.45),
        )

        banner.fill.solid()
        banner.fill.fore_color.rgb = RGBColor(
            0,
            70,
            140,
        )
        banner.line.fill.background()

        title_box = ppt_slide.shapes.title

        title = "ResearchX"

        body = None

        if slide_number != 1:
            body = ppt_slide.placeholders[1].text_frame
            body.clear()

        bullet_mode = False

        # ----------------------------------------
        # Parse AI Output
        # ----------------------------------------
        for line in slide.splitlines():

            line = line.strip()

            if not line:
                continue

            if line.startswith("**Slide Title:**"):

                title = (
                    line.replace(
                        "**Slide Title:**",
                        "",
                    )
                    .replace("**", "")
                    .strip()
                )

                continue

            if "Bullet Points" in line:
                bullet_mode = True
                continue

            if "Speaker Notes" in line:
                bullet_mode = False
                continue

            if bullet_mode and body is not None:

                line = (
                    line.replace("-", "")
                    .replace("*", "")
                    .replace("•", "")
                    .strip()
                )

                if not line:
                    continue

                p = body.add_paragraph()

                p.text = line

                p.level = 0

                p.font.size = Pt(22)

                p.font.color.rgb = RGBColor(
                    40,
                    40,
                    40,
                )

        # ----------------------------------------
        # Title
        # ----------------------------------------
        title_box.text = title

        title_para = title_box.text_frame.paragraphs[0]

        title_para.font.size = Pt(30)

        title_para.font.bold = True

        title_para.font.color.rgb = RGBColor(
            0,
            70,
            140,
        )

        # ----------------------------------------
        # First Slide Subtitle
        # ----------------------------------------
        if slide_number == 1:

            subtitle = ppt_slide.placeholders[1]

            subtitle.text = (
                "AI Multi-Agent Research Assistant\n\n"
                "Automatically Generated Presentation\n\n"
                "ResearchX • 2026"
            )

            sub = subtitle.text_frame.paragraphs[0]

            sub.alignment = PP_ALIGN.CENTER

            sub.font.size = Pt(22)

            sub.font.bold = True

            sub.font.color.rgb = RGBColor(
                70,
                70,
                70,
            )

        # ----------------------------------------
        # Footer
        # ----------------------------------------
        footer = ppt_slide.shapes.add_textbox(
            Inches(0.3),
            Inches(7.0),
            Inches(4),
            Inches(0.3),
        )

        fp = footer.text_frame.paragraphs[0]

        fp.text = "ResearchX"

        fp.font.size = Pt(10)

        fp.font.color.rgb = RGBColor(
            120,
            120,
            120,
        )

        # ----------------------------------------
        # Slide Number
        # ----------------------------------------
        number = ppt_slide.shapes.add_textbox(
            Inches(12.5),
            Inches(7.0),
            Inches(0.5),
            Inches(0.3),
        )

        np = number.text_frame.paragraphs[0]

        np.text = str(slide_number)

        np.alignment = PP_ALIGN.RIGHT

        np.font.size = Pt(10)

        np.font.color.rgb = RGBColor(
            120,
            120,
            120,
        )

        slide_number += 1

    prs.save(output_path)

    return output_path