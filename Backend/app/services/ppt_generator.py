import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def generate_ieee_presentation(paper_title: str, slides_content: list[dict]) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    NAVY = RGBColor(14, 43, 92)
    GRAY = RGBColor(40, 40, 40)

    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    title.text = paper_title[:75]
    title.text_frame.paragraphs[0].font.color.rgb = NAVY
    title.text_frame.paragraphs[0].font.bold = True
    title_slide.placeholders[1].text = "Technical Presentation\nResearchX IEEE Standard Deck"

    # Content Slides
    for slide_data in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        s_title = slide.shapes.title
        s_title.text = slide_data.get("heading", "Analysis")
        s_title.text_frame.paragraphs[0].font.color.rgb = NAVY
        s_title.text_frame.paragraphs[0].font.bold = True

        tf = slide.shapes.placeholders[1].text_frame
        tf.word_wrap = True
        tf.text = ""
        for bullet in slide_data.get("points", []):
            p = tf.add_paragraph()
            p.text = f"•  {bullet}"
            p.font.size = Pt(17)
            p.font.color.rgb = GRAY

    os.makedirs("uploads/generated_ppts", exist_ok=True)
    clean_name = "".join(c for c in paper_title[:20] if c.isalnum() or c in (' ', '_')).strip().replace(' ', '_')
    file_path = f"uploads/generated_ppts/{clean_name or 'paper'}.pptx"
    prs.save(file_path)
    return file_path